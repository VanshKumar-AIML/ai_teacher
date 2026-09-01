import PyPDF2
import docx
from pptx import Presentation
import io
import os

def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_txt(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore")

def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == ".docx":
        return extract_text_from_docx(file_bytes)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_bytes)
    elif ext == ".txt":
        return extract_text_from_txt(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")